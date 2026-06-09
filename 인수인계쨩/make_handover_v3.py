"""
성과관리 인수인계 일정안내 PPT v3
U+design.md 디자인 가이드 완전 반영
- 배경: #f4f3ef (Warm Gray)
- 카드: #ffffff (White)
- 브랜드 액센트: #b90064 (U+ Magenta)
- 텍스트: #1c1b1b (Carbon) / #5c5a54 (Onvariant)
- 테두리: #e9e8e3
- 총합계 행: #1c1b1b 배경 + 흰 텍스트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── U+ 색상 팔레트 ────────────────────────────
BG          = RGBColor(0xF4, 0xF3, 0xEF)   # brand-bg (Warm Gray)
SURFACE     = RGBColor(0xFF, 0xFF, 0xFF)   # 카드/콘텐츠 박스
SURFACE_SUB = RGBColor(0xFA, 0xFA, 0xF9)   # 섹션 내부 그룹 배경
DIVIDER     = RGBColor(0xE9, 0xE8, 0xE3)   # 기본 테두리/디바이더
DIVIDER_HI  = RGBColor(0xD4, 0xD3, 0xCC)   # 강조 경계선
CARBON      = RGBColor(0x1C, 0x1B, 0x1B)   # 기본 텍스트
ONVARIANT   = RGBColor(0x5C, 0x5A, 0x54)   # 보조 텍스트
OUTLINE     = RGBColor(0x8A, 0x88, 0x80)   # 비활성/플레이스홀더
MAGENTA     = RGBColor(0xB9, 0x00, 0x64)   # U+ 브랜드 마젠타
AMBER       = RGBColor(0xD9, 0x77, 0x06)   # 주의/경보
AMBER_BG    = RGBColor(0xFD, 0xF6, 0xEE)   # 주의 배경 (~10%)
TEAL        = RGBColor(0x0D, 0x94, 0x88)   # 특수 처리
TEAL_BG     = RGBColor(0xED, 0xF6, 0xF4)   # 특수 처리 배경 (~6%)
TBL_HDR     = RGBColor(0xF5, 0xF0, 0xE8)   # 테이블 헤더 배경
TOTAL_BG    = RGBColor(0x1C, 0x1B, 0x1B)   # 총합계 행 배경 (Carbon)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# ── 네비 전용 색 ──────────────────────────────
NAV_BG      = RGBColor(0x18, 0x17, 0x15)   # 네비 배경 (더 깊은 다크)
NAV_ACTIVE  = RGBColor(0x26, 0x24, 0x21)   # 활성 아이템 배경
NAV_DIM     = RGBColor(0x3E, 0x3C, 0x38)   # 비활성 아이콘 배경
NAV_TEXT    = RGBColor(0xC8, 0xC6, 0xC0)   # 비활성 텍스트
NAV_SUB     = RGBColor(0x6E, 0x6C, 0x66)   # 비활성 서브텍스트
NAV_DIV     = RGBColor(0x2C, 0x2A, 0x27)   # 네비 내부 구분선

# ── 슬라이드 크기 16:9 ────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

W = Inches(13.33)
H = Inches(7.5)

NAV_W     = Inches(2.45)
CONTENT_L = NAV_W
CONTENT_W = W - NAV_W
HDR_H     = Inches(0.82)

SECTIONS = [
    {"id": "pre",   "label": "데이터 생성 전", "sub": "월초 ~ 11일",      "num": "01"},
    {"id": "batch", "label": "정기 배치",       "sub": "매월 11일 (자동)", "num": "02"},
    {"id": "post",  "label": "데이터 생성 후",  "sub": "11일 이후",        "num": "03"},
]

# ══════════════════════════════════════════════
# 헬퍼
# ══════════════════════════════════════════════

def bg(slide, color=BG):
    b = slide.background; b.fill.solid(); b.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line=None, lw=Pt(0.75)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else: sh.line.fill.background()
    return sh

def rrect(slide, l, t, w, h, fill=None, line=None, adj=0.07, lw=Pt(0.75)):
    sh = slide.shapes.add_shape(5, l, t, w, h)
    sh.adjustments[0] = adj
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else: sh.line.fill.background()
    return sh

def oval(slide, l, t, d, fill):
    sh = slide.shapes.add_shape(9, l, t, d, d)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=11, bold=False, color=CARBON,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Arial"
    return tb


# ══════════════════════════════════════════════
# 공통 레이아웃: 네비 바 + 콘텐츠 헤더
# ══════════════════════════════════════════════

def draw_layout(slide, active_id, title, sub=""):
    bg(slide)

    # ── 네비 배경 ────────────────────────────
    rect(slide, 0, 0, NAV_W, H, NAV_BG)

    # 네비 상단 로고/브랜드 영역
    rect(slide, 0, 0, NAV_W, Inches(0.82), RGBColor(0x10, 0x0F, 0x0D))
    # 마젠타 포인트 라인
    rect(slide, 0, 0, Inches(0.04), Inches(0.82), MAGENTA)
    txt(slide, "LG U+  성과관리", Inches(0.14), Inches(0.14),
        NAV_W - Inches(0.2), Inches(0.3),
        size=11, bold=True, color=WHITE)
    txt(slide, "인수인계 일정 안내", Inches(0.14), Inches(0.46),
        NAV_W - Inches(0.2), Inches(0.22),
        size=8, color=RGBColor(0x8A, 0x88, 0x80))

    # ── 섹션 아이템 ──────────────────────────
    SEC_H = Inches(1.3)
    sy = Inches(0.92)

    for sec in SECTIONS:
        is_active = sec["id"] == active_id

        if is_active:
            # 활성: 마젠타 사이드 바 + 약간 밝은 배경
            rect(slide, 0, sy, Inches(0.04), SEC_H, MAGENTA)
            rect(slide, Inches(0.04), sy, NAV_W - Inches(0.04), SEC_H, NAV_ACTIVE)

            # 번호 뱃지 (마젠타 원)
            d = Inches(0.44)
            oval(slide, Inches(0.18), sy + Inches(0.28), d, MAGENTA)
            txt(slide, sec["num"],
                Inches(0.18), sy + Inches(0.28), d, d,
                size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

            txt(slide, sec["label"],
                Inches(0.74), sy + Inches(0.2),
                NAV_W - Inches(0.84), Inches(0.36),
                size=13, bold=True, color=WHITE)
            txt(slide, sec["sub"],
                Inches(0.74), sy + Inches(0.58),
                NAV_W - Inches(0.84), Inches(0.22),
                size=9, color=RGBColor(0xE8, 0xA0, 0xC4))  # 연마젠타
        else:
            d = Inches(0.36)
            oval(slide, Inches(0.22), sy + Inches(0.3), d, NAV_DIM)
            txt(slide, sec["num"],
                Inches(0.22), sy + Inches(0.3), d, d,
                size=10, bold=True, color=OUTLINE, align=PP_ALIGN.CENTER)
            txt(slide, sec["label"],
                Inches(0.70), sy + Inches(0.24),
                NAV_W - Inches(0.8), Inches(0.32),
                size=11, color=NAV_TEXT)
            txt(slide, sec["sub"],
                Inches(0.70), sy + Inches(0.58),
                NAV_W - Inches(0.8), Inches(0.22),
                size=8, color=NAV_SUB)

        # 구분선
        if sec != SECTIONS[-1]:
            rect(slide, Inches(0.1), sy + SEC_H - Inches(0.008),
                 NAV_W - Inches(0.2), Inches(0.008), NAV_DIV)

        sy += SEC_H

    # 네비 하단: 비정기 작업
    irr_y = H - Inches(1.05)
    rect(slide, Inches(0.1), irr_y, NAV_W - Inches(0.2), Inches(0.008), NAV_DIV)
    txt(slide, "비정기 작업", Inches(0.18), irr_y + Inches(0.1),
        NAV_W - Inches(0.28), Inches(0.26),
        size=9, bold=True, color=RGBColor(0x6E, 0x6C, 0x66))
    txt(slide, "필요 시 수행", Inches(0.18), irr_y + Inches(0.36),
        NAV_W - Inches(0.28), Inches(0.2),
        size=8, color=NAV_SUB)

    # ── 콘텐츠 헤더 ──────────────────────────
    # 흰 카드 영역 전체
    rect(slide, CONTENT_L, 0, CONTENT_W, H, BG)

    # 헤더 바: 흰 배경 + 하단 마젠타 라인
    rect(slide, CONTENT_L, 0, CONTENT_W, HDR_H, SURFACE)
    rect(slide, CONTENT_L, HDR_H - Inches(0.04), CONTENT_W, Inches(0.04), MAGENTA)

    txt(slide, title,
        CONTENT_L + Inches(0.35), Inches(0.1),
        CONTENT_W - Inches(0.5), Inches(0.44),
        size=20, bold=True, color=CARBON)
    if sub:
        txt(slide, sub,
            CONTENT_L + Inches(0.35), Inches(0.54),
            CONTENT_W - Inches(0.5), Inches(0.24),
            size=10, color=ONVARIANT)

    return HDR_H + Inches(0.22)


# ══════════════════════════════════════════════
# 표 그리기 헬퍼
# ══════════════════════════════════════════════

def draw_table(slide, start_y, rows_data, show_total=True):
    HDR_Y = start_y
    HDR_H_ROW = Inches(0.34)

    CL = [
        CONTENT_L + Inches(0.22),
        CONTENT_L + Inches(2.0),
        CONTENT_L + Inches(5.65),
        CONTENT_L + Inches(9.85),
    ]
    CW = [Inches(1.73), Inches(3.6), Inches(4.15), Inches(1.25)]
    HEADERS = ["구분", "작업 항목", "세부 내용", "소요시간"]

    # ── 헤더 행 ──────────────────────────────
    # 헤더 배경 (수평선만, 수직선 없음 — U+design 가이드 준수)
    total_w = sum(CW) + Inches(0.05)
    rect(slide, CL[0], HDR_Y, total_w, HDR_H_ROW, TBL_HDR)
    # 하단 수평선만
    rect(slide, CL[0], HDR_Y + HDR_H_ROW - Inches(0.01),
         total_w, Inches(0.01), DIVIDER)

    for cl, cw, hd in zip(CL, CW, HEADERS):
        txt(slide, hd, cl + Inches(0.08), HDR_Y + Inches(0.07),
            cw - Inches(0.12), HDR_H_ROW - Inches(0.1),
            size=9, bold=True, color=ONVARIANT,
            align=PP_ALIGN.CENTER if hd == "소요시간" else PP_ALIGN.LEFT)

    ROW_H = Inches(0.72)
    ry = HDR_Y + HDR_H_ROW
    prev_cat = None

    for row in rows_data:
        cat, task, detail, row_bg = row[0], row[1], row[2], row[3]
        duration = row[4] if len(row) > 4 else ""

        # 행 배경 — 흰 카드 or 연한 서브
        fill = SURFACE_SUB if row_bg == "sub" else SURFACE
        rect(slide, CL[0], ry, total_w, ROW_H, fill)
        # 하단 수평선만 (수직선 없음)
        rect(slide, CL[0], ry + ROW_H - Inches(0.008),
             total_w, Inches(0.008), DIVIDER)

        # 구분 셀: 같은 구분이면 비움
        if cat and cat != prev_cat:
            # 마젠타 틴트 배경 캡슐
            rrect(slide, CL[0] + Inches(0.08), ry + Inches(0.14),
                  CW[0] - Inches(0.16), ROW_H - Inches(0.28),
                  fill=MAGENTA, adj=0.12)
            txt(slide, cat,
                CL[0] + Inches(0.08), ry + Inches(0.14),
                CW[0] - Inches(0.16), ROW_H - Inches(0.28),
                size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            prev_cat = cat
        # 같은 cat이거나 None이면 셀 비움 (화살표 없음)

        # 작업 항목 — 번호 부분 마젠타
        task_parts = task.split(" ", 1)
        if task_parts[0] in ["①","②","③","④","⑤","⑥","⑦"]:
            tb = slide.shapes.add_textbox(
                CL[1] + Inches(0.1), ry + Inches(0.12),
                CW[1] - Inches(0.15), ROW_H - Inches(0.18))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r1 = p.add_run(); r1.text = task_parts[0] + "  "
            r1.font.size = Pt(10); r1.font.bold = True
            r1.font.color.rgb = MAGENTA; r1.font.name = "Arial"
            r2 = p.add_run(); r2.text = task_parts[1] if len(task_parts) > 1 else ""
            r2.font.size = Pt(10); r2.font.bold = True
            r2.font.color.rgb = CARBON; r2.font.name = "Arial"
        else:
            txt(slide, task,
                CL[1] + Inches(0.1), ry + Inches(0.12),
                CW[1] - Inches(0.15), ROW_H - Inches(0.18),
                size=10, bold=True, color=CARBON)

        # 세부 내용 — Caption 스타일
        txt(slide, detail,
            CL[2] + Inches(0.1), ry + Inches(0.08),
            CW[2] - Inches(0.15), ROW_H - Inches(0.1),
            size=9, color=ONVARIANT)

        # 소요시간 — 오른쪽 정렬, tabular 느낌
        if duration:
            txt(slide, duration,
                CL[3] + Inches(0.05), ry + Inches(0.2),
                CW[3] - Inches(0.08), ROW_H - Inches(0.3),
                size=11, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)

        ry += ROW_H

    # ── 총합계 행 — Carbon(#1c1b1b) 배경 + 흰 텍스트 ──
    if show_total:
        rect(slide, CL[0], ry, total_w, Inches(0.42), TOTAL_BG)
        txt(slide, "소요시간 합계",
            CL[0] + Inches(0.15), ry + Inches(0.1),
            Inches(6), Inches(0.26),
            size=9, bold=True, color=WHITE)

        durations = [r[4] for r in rows_data if len(r) > 4 and r[4]]
        if durations:
            total_md = sum(float(d.replace("MD","").strip()) for d in durations)
            total_label = f"{total_md:g} MD"
            txt(slide, total_label,
                CL[3] + Inches(0.05), ry + Inches(0.08),
                CW[3] - Inches(0.08), Inches(0.28),
                size=12, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)

    return ry + (Inches(0.42) if show_total else 0)


# ══════════════════════════════════════════════
# SLIDE 1 — 데이터 생성 전 (통합)
# ══════════════════════════════════════════════
sl1 = prs.slides.add_slide(BLANK)
cy = draw_layout(sl1, "pre",
    "데이터 생성 전 — 수기테이블 생성 및 실적 조정 데이터 준비",
    "월초 ~ 매월 11일 | 마스터·분류 정보 확정 + 조정 대상 실적 확인 + PL 이익률 수기테이블 생성")

rows_1 = [
    ("KA/GA\n고객 리스트",
     "① KA/GA 고객 리스트 확정",
     "현행 KA·GA 고객 목록 검토 및 변경사항 반영\nKAM TF 협의 기반 최종 확정 (KAM TF 협의 및 검증 작업 필요)",
     "white", "0.5 MD"),
    ("KA/GA\n고객 리스트",
     "② 고객 그룹핑 정보 확정",
     "복수 법인을 단일 그룹으로 묶는 매핑 정보 검토·갱신\nKAM TF 협의 기반 최종 확정 (KAM TF 협의 및 검증 작업 필요)",
     "white", "0.5 MD"),
    ("고객군\n(Lv1~4)",
     "③ 고객군 분류 변경사항 점검 및 반영",
     "각 고객별 고객군 분류 변경사항 확인 및 반영\n기존 데이터 중 변경 요청사항 반영, 공정위 대상 대기업집단 업데이트 등\nKAM TF 협의 및 검증 작업 필요",
     "sub", "1 MD"),
    ("매출/이익\n조정 점검",
     "④ DW 데이터 생성 · 적재 스케줄 점검",
     "기업DW 데이터 생성 완료 여부 확인\n데이터레이크(DL) 적재 스케줄 이상 없음 점검",
     "white", "0.1 MD"),
    ("매출/이익\n조정 점검",
     "⑤ 주요 조정 대상 실적 확인",
     "네이버 청구매출, 에스원·KGM·토요타 청구매출, MVNO 정산매출 확인 등\n조정 대상 고객 실적 추출·확인 (기업DW, 파트너플러스 시스템 내 실적 확인 필요)",
     "white", "0.5 MD"),
    ("PL 이익률\n수기테이블",
     "⑥ 부문 PL 이익률 마감 스케줄 확인",
     "기획팀 PL 자료 제공 일정 확인 (통상 6영업일 전후)",
     "sub", "0.1 MD"),
    ("PL 이익률\n수기테이블",
     "⑦ 7개 수기테이블 생성/검증 후 운영 업로드",
     "월별 업데이트 필요한 7개 수기테이블 생성/검증 후 레이크 업로드",
     "sub", "2 MD"),
]

end_y = draw_table(sl1, cy, rows_1, show_total=True)

txt(sl1, "※ KAM TF 협의 일정에 따라 확정 시점 유동   ※ PL 이익률 수기테이블 소요시간은 수행 후 기입 예정",
    CONTENT_L + Inches(0.22), end_y + Inches(0.1),
    CONTENT_W - Inches(0.3), Inches(0.22),
    size=8, italic=True, color=OUTLINE)


# ══════════════════════════════════════════════
# SLIDE 2 — 정기 배치
# ══════════════════════════════════════════════
sl2 = prs.slides.add_slide(BLANK)
cy2 = draw_layout(sl2, "batch",
    "정기 배치",
    "매월 11일 자동 배치 | 수기데이터 완성 여부에 따라 이후 작업 흐름 분기")

CX = CONTENT_L + Inches(0.22)
CW2 = CONTENT_W - Inches(0.44)

# ── 상단: 정기 배치 설명 박스 ─────────────────
rrect(sl2, CX, cy2, CW2, Inches(1.0),
      fill=SURFACE, line=DIVIDER, adj=0.05)
# 왼쪽 마젠타 포인트 바
rect(sl2, CX, cy2, Inches(0.04), Inches(1.0), MAGENTA)
txt(sl2, "정기 배치 개요",
    CX + Inches(0.16), cy2 + Inches(0.1),
    CW2 - Inches(0.22), Inches(0.28),
    size=11, bold=True, color=CARBON)
txt(sl2,
    "정기 배치는 매월 11일 자동 배치되나, 정기 배치 이전에 작업되어야 하는 수기데이터 작업 미완성 시 "
    "정기 배치 이후 수기데이터 생성 → 반영 → 재배치 작업이 필요함",
    CX + Inches(0.16), cy2 + Inches(0.42),
    CW2 - Inches(0.22), Inches(0.48),
    size=10, color=ONVARIANT)

cy2 += Inches(1.0) + Inches(0.18)

# ── 분기 흐름도 ──────────────────────────────
# 공통 시작: 매월 11일 배치
NODE_W = Inches(2.4)
NODE_H = Inches(0.72)
MID_X  = CONTENT_L + CONTENT_W / 2 - NODE_W / 2  # 가운데 정렬

# 시작 노드 (마젠타)
rrect(sl2, MID_X, cy2, NODE_W, NODE_H,
      fill=MAGENTA, line=None, adj=0.08)
txt(sl2, "매월 11일  정기 배치 실행",
    MID_X, cy2, NODE_W, NODE_H,
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

cy2 += NODE_H + Inches(0.1)

# 분기 화살표 + 라벨
FORK_H = Inches(0.28)
# 세로 줄기
rect(sl2, MID_X + NODE_W/2 - Inches(0.01), cy2,
     Inches(0.02), FORK_H, ONVARIANT)
cy2 += FORK_H

# 다이아몬드 분기 라벨
DIAG_W = Inches(3.2)
DIAG_H = Inches(0.52)
diag_x = CONTENT_L + CONTENT_W / 2 - DIAG_W / 2
rrect(sl2, diag_x, cy2, DIAG_W, DIAG_H,
      fill=TBL_HDR, line=DIVIDER_HI, adj=0.08)
txt(sl2, "배치 전 수기데이터 작업 완료 여부?",
    diag_x, cy2, DIAG_W, DIAG_H,
    size=9, bold=True, color=CARBON, align=PP_ALIGN.CENTER)

cy2 += DIAG_H + Inches(0.1)

# 좌(미완성) / 우(정상) 분기선
BRANCH_Y  = cy2
LEFT_CX   = CONTENT_L + Inches(2.0)
RIGHT_CX  = CONTENT_L + CONTENT_W - Inches(2.0)

# 수평 분기선
rect(sl2, LEFT_CX, BRANCH_Y, RIGHT_CX - LEFT_CX, Inches(0.02), ONVARIANT)
# 좌 수직
rect(sl2, LEFT_CX - Inches(0.01), BRANCH_Y, Inches(0.02), Inches(0.28), ONVARIANT)
# 우 수직
rect(sl2, RIGHT_CX - Inches(0.01), BRANCH_Y, Inches(0.02), Inches(0.28), ONVARIANT)

# 분기 라벨
txt(sl2, "미완성", LEFT_CX - Inches(0.3), BRANCH_Y - Inches(0.22),
    Inches(0.9), Inches(0.22), size=8, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
txt(sl2, "완료", RIGHT_CX - Inches(0.3), BRANCH_Y - Inches(0.22),
    Inches(0.9), Inches(0.22), size=8, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

cy2 += Inches(0.28)

# ── 왼쪽 분기: 미완성 (Amber 계열) ─────────────
BOX_W = Inches(4.0)
BOX_H = Inches(1.55)
LEFT_BX = LEFT_CX - BOX_W / 2

rrect(sl2, LEFT_BX, cy2, BOX_W, BOX_H,
      fill=AMBER_BG, line=AMBER, adj=0.06, lw=Pt(1.0))
rect(sl2, LEFT_BX, cy2, Inches(0.04), BOX_H, AMBER)

txt(sl2, "수기데이터 미완성 시",
    LEFT_BX + Inches(0.16), cy2 + Inches(0.1),
    BOX_W - Inches(0.22), Inches(0.26),
    size=10, bold=True, color=AMBER)

steps_left = [
    "① 수기데이터 생성 완료",
    "② 시스템 반영",
    "③ 재배치 실행",
]
for i, s in enumerate(steps_left):
    sy_item = cy2 + Inches(0.42) + i * Inches(0.34)
    oval(sl2, LEFT_BX + Inches(0.16), sy_item + Inches(0.05),
         Inches(0.22), AMBER)
    txt(sl2, str(i+1),
        LEFT_BX + Inches(0.16), sy_item + Inches(0.05),
        Inches(0.22), Inches(0.22),
        size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl2, s.split(" ", 1)[1],
        LEFT_BX + Inches(0.46), sy_item,
        BOX_W - Inches(0.56), Inches(0.3),
        size=9, color=CARBON)

# 재배치 후 → 생성 후 작업으로 화살표
arr_y = cy2 + BOX_H + Inches(0.06)
rect(sl2, LEFT_CX - Inches(0.01), arr_y, Inches(0.02), Inches(0.28), ONVARIANT)
txt(sl2, "재배치 완료 후 →",
    LEFT_BX, arr_y + Inches(0.04),
    BOX_W, Inches(0.22),
    size=8, italic=True, color=ONVARIANT, align=PP_ALIGN.CENTER)

# ── 오른쪽 분기: 정상 (Teal 계열) ──────────────
RIGHT_BX = RIGHT_CX - BOX_W / 2

rrect(sl2, RIGHT_BX, cy2, BOX_W, BOX_H,
      fill=TEAL_BG, line=RGBColor(0xA8,0xD9,0xD0), adj=0.06, lw=Pt(1.0))
rect(sl2, RIGHT_BX, cy2, Inches(0.04), BOX_H, TEAL)

txt(sl2, "정상 배치 완료 시",
    RIGHT_BX + Inches(0.16), cy2 + Inches(0.1),
    BOX_W - Inches(0.22), Inches(0.26),
    size=10, bold=True, color=TEAL)
txt(sl2, "데이터 생성 이후 작업으로 바로 이동",
    RIGHT_BX + Inches(0.16), cy2 + Inches(0.42),
    BOX_W - Inches(0.22), Inches(0.28),
    size=9, color=CARBON)

# 화살표
arr_y2 = cy2 + BOX_H + Inches(0.06)
rect(sl2, RIGHT_CX - Inches(0.01), arr_y2, Inches(0.02), Inches(0.28), ONVARIANT)
txt(sl2, "→  데이터 생성 후 작업",
    RIGHT_BX, arr_y2 + Inches(0.04),
    BOX_W, Inches(0.22),
    size=8, italic=True, color=ONVARIANT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# SLIDE 3 — 데이터 생성 후 (표 레이아웃, 1페이지와 동일)
# ══════════════════════════════════════════════
sl3 = prs.slides.add_slide(BLANK)
cy3 = draw_layout(sl3, "post",
    "데이터 생성 후",
    "매월 11일 이후 | KA 고객 성과지표 추출 → 점검 → 조정 → 리뷰 → 최종 확인")

rows_3 = [
    # (구분, 작업항목, 세부내용, 행배경, 소요시간)
    ("성과지표\n추출·점검",
     "① KA 고객 성과지표 추출",
     "KA 고객 중심으로 매출/이익률 성과지표 추출",
     "white", "0.2 MD"),
    ("성과지표\n추출·점검",
     "② 성과지표 이상여부 체크",
     "고객의 손익별 매출/이익 산출 결과 점검 후\n초기 이상여부·예외처리 필요사항 도출",
     "white", "1 MD"),
    ("성과지표\n추출·점검",
     "③ 조정 작업 수행",
     "초기 이상여부·예외처리 조정 작업 수행",
     "white", "0.5 MD"),
    ("성과지표\n추출·점검",
     "④ 성과지표 재추출",
     "조정 작업 반영 후 성과지표 재추출",
     "white", "0.2 MD"),
    ("성과지표\n추출·점검",
     "⑤ KA 고객 성과지표 공유",
     "KAM TF, 담당 내 등 공유",
     "white", "0.1 MD"),
    ("리뷰",
     "⑥ KAM TF 성과지표 리뷰",
     "KAM TF와 KA 성과지표 리뷰",
     "sub", "1 MD"),
    ("리뷰",
     "⑦ 기획담당 성과지표 리뷰",
     "기획담당 KA 성과지표 리뷰",
     "sub", "1 MD"),
    ("리뷰",
     "⑧ 조정 사항 반영",
     "리뷰 결과에 따른 조정 사항 반영 작업 수행",
     "sub", "0.5 MD"),
    ("리뷰",
     "⑨ 최종 조정결과 확인",
     "리뷰 반영 후 최종 산출 결과 확인 및 이상 없음 검증",
     "sub", "0.1 MD"),
    # 병렬 작업 — 구분 셀에 표시
    ("병렬 수행\n(⑨ 이후)",
     "⑩ 경영성과리뷰 양식 작성",
     "Claude 활용 양식 초안 생성 → 검토 및 확정",
     "white", "1 MD"),
    ("병렬 수행\n(⑨ 이후)",
     "⑪ 전체 데이터 및 KA 기타지표 정상 적재 확인",
     "전체 고객 데이터 적재 완료 여부 및 KA 기타지표 정합성 확인",
     "white", "1 MD"),
]

end_y3 = draw_table(sl3, cy3, rows_3, show_total=True)

# 병렬 작업 안내 주석
txt(sl3, "※ ⑩⑪은 ⑨ 최종 조정결과 확인 이후 동시(병렬) 수행 — 합산 소요시간에서 중복 산정됨",
    CONTENT_L + Inches(0.22), end_y3 + Inches(0.1),
    CONTENT_W - Inches(0.3), Inches(0.22),
    size=8, italic=True, color=OUTLINE)


# ══════════════════════════════════════════════
# 저장
# ══════════════════════════════════════════════
out = "/Users/konslie/Desktop/ClaudeCode/인수인계쨩/성과관리_인수인계_일정안내.pptx"
prs.save(out)
print(f"저장 완료: {out}  (총 {len(prs.slides)}장)")
