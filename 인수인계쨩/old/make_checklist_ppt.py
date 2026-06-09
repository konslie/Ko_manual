"""
성과관리 월별 작업 체크리스트 PPT 생성
index6.html 디자인 언어 기반: 크림 배경, 오렌지 강조(#CC785C), Pretendard 스타일
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ──────────────────────────────
BG       = RGBColor(0xFA, 0xF9, 0xF5)   # #faf9f5 크림
CARD     = RGBColor(0xEF, 0xE9, 0xDE)   # #efe9de 카드
MUTED    = RGBColor(0xF5, 0xF0, 0xE8)   # #f5f0e8 연한 배경
BORDER   = RGBColor(0xC8, 0xBF, 0xB4)   # #c8bfb4 테두리
ORANGE   = RGBColor(0xCC, 0x78, 0x5C)   # #cc785c 주강조
ORANGE_D = RGBColor(0xA9, 0x58, 0x3E)   # #a9583e 어두운 오렌지
TEXT     = RGBColor(0x14, 0x14, 0x13)   # #141413 본문
SUB      = RGBColor(0x6C, 0x6A, 0x64)   # #6c6a64 보조
BODY     = RGBColor(0x3D, 0x3D, 0x3A)   # #3d3d3a
GREEN    = RGBColor(0x2E, 0x7A, 0x6A)   # #2e7a6a
GREEN_BG = RGBColor(0xD4, 0xED, 0xE8)   # #d4ede8
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x3D, 0x3D, 0x3A)   # #3d3d3a 어두운
RED_BG   = RGBColor(0xFD, 0xE8, 0xE0)
RED      = RGBColor(0xB9, 0x4A, 0x2D)

# ── 슬라이드 크기 (16:9 와이드) ──────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # 완전 빈 레이아웃

# ══════════════════════════════════════════════
# 헬퍼 함수
# ══════════════════════════════════════════════

def slide_bg(slide, color):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    """사각형 도형 추가 → shape 반환"""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, l, t, w, h)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, l, t, w, h, fill_color=None, line_color=None, radius=0.08):
    """둥근 사각형 (MSO_SHAPE 5 = ROUNDED_RECTANGLE)"""
    shape = slide.shapes.add_shape(5, l, t, w, h)
    shape.adjustments[0] = radius
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, l, t, d, fill_color, line_color=None):
    """원형"""
    shape = slide.shapes.add_shape(9, l, t, d, d)  # 9 = OVAL
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, italic=False,
             word_wrap=True, valign=None):
    """텍스트박스 추가 → shape 반환"""
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"          # PPT 기본 (Pretendard 미설치 환경 대비)
    return txBox


def add_text_multiline(slide, lines, l, t, w, h,
                       size=12, bold=False, color=TEXT,
                       align=PP_ALIGN.LEFT, line_spacing=None):
    """여러 줄 텍스트박스 (list of (text, bold, color, size) or str)"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, c, s = item, bold, color, size
        else:
            txt = item[0]
            b   = item[1] if len(item) > 1 else bold
            c   = item[2] if len(item) > 2 else color
            s   = item[3] if len(item) > 3 else size

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = align
        if line_spacing:
            from pptx.util import Pt as UPt
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(s)
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = "Arial"
    return txBox


def add_circle_label(slide, cx, cy, r, num_text, fill=ORANGE, text_color=WHITE):
    """원형 숫자 레이블"""
    l = cx - r
    t = cy - r
    d = r * 2
    add_circle(slide, l, t, d, fill)
    add_text(slide, num_text, l, t, d, d,
             size=11, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def orange_header(slide, title, subtitle="", num_label=""):
    """오렌지 그라디언트 헤더 바"""
    bar = add_rect(slide, 0, 0, W, Inches(1.1), fill_color=ORANGE)
    if num_label:
        add_text(slide, num_label, Inches(0.4), Inches(0.08), Inches(5), Inches(0.3),
                 size=9, bold=True, color=RGBColor(0xFF,0xCC,0xAA))
    add_text(slide, title, Inches(0.4), Inches(0.3) if num_label else Inches(0.2),
             Inches(12), Inches(0.55),
             size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.82), Inches(12), Inches(0.3),
                 size=11, color=RGBColor(0xFF,0xE8,0xD8))


def check_row(slide, l, t, w, text, sub="", done=False):
    """체크 아이템 행"""
    h = Inches(0.42) if not sub else Inches(0.55)
    bg = GREEN_BG if done else MUTED
    add_rounded_rect(slide, l, t, w, h, fill_color=bg, line_color=BORDER, radius=0.06)
    # 원형 체크
    c_fill = GREEN if done else WHITE
    add_circle(slide, l + Inches(0.1), t + Inches(0.1), Inches(0.24), c_fill, line_color=BORDER)
    if done:
        add_text(slide, "✓", l + Inches(0.1), t + Inches(0.08), Inches(0.24), Inches(0.26),
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 텍스트
    lines = [(text, True, TEXT if not done else SUB, 11)]
    if sub:
        lines.append((sub, False, SUB, 9))
    add_text_multiline(slide, lines, l + Inches(0.42), t + Inches(0.06), w - Inches(0.52), h - Inches(0.1))
    return h


def pill_badge(slide, l, t, text, fill=CARD, text_color=ORANGE, size=9):
    """뱃지 (둥근 사각형)"""
    tw = Inches(len(text) * 0.11 + 0.3)
    add_rounded_rect(slide, l, t, tw, Inches(0.26), fill_color=fill, line_color=None, radius=0.3)
    add_text(slide, text, l, t, tw, Inches(0.26),
             size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return tw


def v_connector(slide, cx, y1, y2):
    """세로 연결선"""
    add_rect(slide, cx - Inches(0.01), y1, Inches(0.02), y2 - y1, fill_color=BORDER)


def sec_pill_label(slide, l, t, num, text):
    """sec-pill 스타일 레이블"""
    d = Inches(0.36)
    add_circle(slide, l, t, d, ORANGE)
    add_text(slide, num, l, t, d, d, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, text, l + d + Inches(0.12), t + Inches(0.02),
             Inches(8), d - Inches(0.04), size=18, bold=True, color=TEXT)


# ══════════════════════════════════════════════
# SLIDE 1 — 표지
# ══════════════════════════════════════════════
sl1 = prs.slides.add_slide(BLANK)
slide_bg(sl1, ORANGE)

# 장식 원
dec = sl1.shapes.add_shape(9, Inches(9.8), Inches(-1.2), Inches(4.5), Inches(4.5))
dec.fill.solid(); dec.fill.fore_color.rgb = RGBColor(0xD4,0x7A,0x55)
dec.line.fill.background()

dec2 = sl1.shapes.add_shape(9, Inches(11.0), Inches(5.2), Inches(3.0), Inches(3.0))
dec2.fill.solid(); dec2.fill.fore_color.rgb = RGBColor(0xD4,0x7A,0x55)
dec2.line.fill.background()

# 태그라인
add_text(sl1, "PERFORMANCE MANAGEMENT · 2026",
         Inches(0.6), Inches(1.2), Inches(8), Inches(0.4),
         size=11, bold=True, color=RGBColor(0xFF,0xCC,0xAA))

# 제목
add_text_multiline(sl1,
    [("성과관리", True, WHITE, 44), ("월별 작업 체크리스트", True, WHITE, 44)],
    Inches(0.6), Inches(1.7), Inches(9), Inches(2.2))

# 부제
add_text(sl1, "데이터 생성 전  →  11일 배치  →  생성 후  →  순서대로 따라하기",
         Inches(0.6), Inches(3.85), Inches(9), Inches(0.4),
         size=14, color=RGBColor(0xFF,0xE8,0xD8))

# 통계 칩 3개
chips = [
    ("3개", "매월 고정 작업"),
    ("2개", "조건부 작업"),
    ("4개", "비정기 작업"),
    ("매월 11일", "배치 기준일"),
]
for i, (val, lbl) in enumerate(chips):
    cx = Inches(0.6 + i * 2.1)
    chip = add_rounded_rect(sl1, cx, Inches(4.5), Inches(1.9), Inches(1.1),
                             fill_color=RGBColor(0xD4,0x7A,0x55), radius=0.1)
    add_text(sl1, val, cx, Inches(4.55), Inches(1.9), Inches(0.55),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl1, lbl, cx, Inches(5.05), Inches(1.9), Inches(0.4),
             size=10, color=RGBColor(0xFF,0xE0,0xCC), align=PP_ALIGN.CENTER)

# 힌트
add_text(sl1, "← → 키 또는 슬라이드 이동으로 체크리스트 확인",
         Inches(0.6), Inches(6.9), Inches(8), Inches(0.35),
         size=10, color=RGBColor(0xFF,0xCC,0xAA))


# ══════════════════════════════════════════════
# SLIDE 2 — 월별 작업 타임라인 개요
# ══════════════════════════════════════════════
sl2 = prs.slides.add_slide(BLANK)
slide_bg(sl2, BG)
orange_header(sl2, "월별 작업 타임라인 — 전체 흐름", "데이터 생성 전 → 배치 → 생성 후 흐름 한눈에 보기", "Overview")

# 타임라인 아이템 정의
items = [
    ("📋", MUTED, BORDER,    "데이터 생성 전",            "에스원·KGM·토요타 보정매출 기입 + 한전 PL 수익률 완성",   "매월 필수",    ORANGE,    CARD),
    ("⚡", ORANGE, None,     "매월 11일 05:00~ — 배치 자동 실행", "담당자 액션 없음, 스케줄러 자동 처리",                 "",           None,      None),
    ("🔁", CARD, ORANGE,     "수기테이블 미완성 시",        "수기테이블 완성 → 데이터 재생성 요청",                    "IF",         DARK,      WHITE),
    ("📊", CARD, ORANGE,     "데이터 생성 완료 시",         "네이버클라우드·한전 수익성 조정 → KA 특이사항 체크",       "ELSE",       ORANGE,    WHITE),
    ("✅", ORANGE, None,     "Eye-Checking",               "Claude 양식 생성 → 항목별 육안 확인",                     "최종 확인",   GREEN,     WHITE),
]

tl_x   = Inches(0.6)
dot_d  = Inches(0.52)
dot_cx = tl_x + dot_d / 2
body_l = tl_x + dot_d + Inches(0.2)
body_w = Inches(11.5)
y_start = Inches(1.35)
row_h   = Inches(0.88)
line_h  = Inches(0.18)

for i, (icon, dot_fill, dot_line, title, sub, badge, badge_fill, badge_txt) in enumerate(items):
    y = y_start + i * (row_h + line_h)

    # 아이콘 원
    dot = sl2.shapes.add_shape(9, tl_x, y, dot_d, dot_d)
    dot.fill.solid(); dot.fill.fore_color.rgb = dot_fill
    if dot_line:
        dot.line.color.rgb = dot_line; dot.line.width = Pt(1.5)
    else:
        dot.line.fill.background()
    add_text(sl2, icon, tl_x, y, dot_d, dot_d, size=18, align=PP_ALIGN.CENTER)

    # 연결선 (마지막 제외)
    if i < len(items) - 1:
        line_y = y + dot_d
        add_rect(sl2, dot_cx - Inches(0.01), line_y, Inches(0.02), line_h, fill_color=BORDER)

    # 본문
    title_lines = [(title, True, TEXT, 14)]
    if badge and badge_fill:
        pass  # 뱃지는 별도
    add_text_multiline(sl2, title_lines, body_l, y + Inches(0.04), body_w - Inches(2.5), Inches(0.32))
    add_text(sl2, sub, body_l, y + Inches(0.34), body_w - Inches(2.5), Inches(0.28), size=10, color=SUB)

    # 뱃지
    if badge and badge_fill:
        bw = Inches(max(len(badge) * 0.12 + 0.25, 0.8))
        b_shape = add_rounded_rect(sl2, Inches(11.8), y + Inches(0.12), bw, Inches(0.28),
                                    fill_color=badge_fill, radius=0.3)
        add_text(sl2, badge, Inches(11.8), y + Inches(0.12), bw, Inches(0.28),
                 size=9, bold=True, color=badge_txt if badge_txt else WHITE, align=PP_ALIGN.CENTER)

# 비정기 별도 박스
bx_y = y_start + len(items) * (row_h + line_h) - Inches(0.1)
add_rounded_rect(sl2, Inches(1.0), bx_y, Inches(11.5), Inches(0.46),
                  fill_color=MUTED, line_color=BORDER, radius=0.05)
add_text(sl2, "비정기 작업 (필요시만)  —  수기테이블 정의 업데이트 · 고객 마스터 정비",
         Inches(1.2), bx_y + Inches(0.08), Inches(11), Inches(0.3), size=10, color=SUB)


# ══════════════════════════════════════════════
# SLIDE 3 — 데이터 생성 전: 수기테이블 정비
# ══════════════════════════════════════════════
sl3 = prs.slides.add_slide(BLANK)
slide_bg(sl3, BG)
orange_header(sl3, "수기테이블 정비 — 11일 배치 이전 완료 필수",
              "에스원·KGM·토요타 보정매출 기입 + 한전 PL 수익률 완성", "01 · 데이터 생성 전")

# 왼쪽 패널: 체크리스트
lp_l = Inches(0.4); lp_t = Inches(1.3); lp_w = Inches(6.0); lp_h = Inches(5.8)
add_rounded_rect(sl3, lp_l, lp_t, lp_w, lp_h, fill_color=CARD, line_color=BORDER, radius=0.04)

add_text(sl3, "📦  보정매출 수기 기입", lp_l + Inches(0.2), lp_t + Inches(0.15), Inches(4.0), Inches(0.35),
         size=13, bold=True, color=TEXT)
pill_badge(sl3, lp_l + Inches(4.3), lp_t + Inches(0.2), "월평균 4~5억 원",
           fill=RED_BG, text_color=RED, size=9)

checks_left = [
    ("에스원  보정매출 → 매출조정 수기테이블 기입", "손익코드: M21 / M2M"),
    ("KGM  보정매출 → 매출조정 수기테이블 기입",   "손익코드: T35 / 커넥티드카"),
    ("토요타  보정매출 → 매출조정 수기테이블 기입", "손익코드: T35 / 커넥티드카"),
    ("한전 PL 수익률  수기테이블 완성",             "기획팀 PL 자료 수령 후 즉시 입력"),
]
cy = lp_t + Inches(0.6)
for txt, sub in checks_left:
    rh = check_row(sl3, lp_l + Inches(0.15), cy, lp_w - Inches(0.3), txt, sub)
    cy += rh + Inches(0.1)

# 타이밍 경고 박스
warn_y = cy + Inches(0.05)
add_rounded_rect(sl3, lp_l + Inches(0.15), warn_y, lp_w - Inches(0.3), Inches(0.7),
                  fill_color=RGBColor(0xFF,0xF5,0xEC), line_color=ORANGE, radius=0.05)
add_text(sl3, "⚡ 타이밍 주의",
         lp_l + Inches(0.3), warn_y + Inches(0.05), Inches(4), Inches(0.25),
         size=9, bold=True, color=ORANGE)
add_text(sl3, "DW 매출·PL 자료는 11일 전후 유동적 — 수령 즉시 기입",
         lp_l + Inches(0.3), warn_y + Inches(0.28), Inches(5.3), Inches(0.3),
         size=9, color=BODY)

# 오른쪽 패널: 확인 방법 A/B/C
rp_l = Inches(6.7); rp_t = Inches(1.3); rp_w = Inches(6.3)
add_text(sl3, "금액 확인 방법", rp_l, rp_t, rp_w, Inches(0.3),
         size=10, bold=True, color=SUB)

options = [
    ("A", ORANGE, WHITE, "기업 DW — 영업사원별 손익 보정매출",        "가장 직접적. DW 접근 권한 필요"),
    ("B", ORANGE, WHITE, "파트너플러스 — 사업자별 정산금액",          "DW 대비 쉬움. 계정 체계 매핑 주의"),
    ("C", MUTED,  SUB,   "담당 영업사원에게 직접 확인",               "A/B 불가 시 최후 수단"),
]
oy = rp_t + Inches(0.4)
for letter, lf, lc, title, desc in options:
    add_rounded_rect(sl3, rp_l, oy, rp_w, Inches(1.1), fill_color=MUTED, line_color=BORDER, radius=0.05)
    # 레터 원
    add_rounded_rect(sl3, rp_l + Inches(0.12), oy + Inches(0.25), Inches(0.52), Inches(0.52),
                      fill_color=lf, radius=0.12)
    add_text(sl3, letter, rp_l + Inches(0.12), oy + Inches(0.25), Inches(0.52), Inches(0.52),
             size=14, bold=True, color=lc, align=PP_ALIGN.CENTER)
    add_text(sl3, title, rp_l + Inches(0.75), oy + Inches(0.15), rp_w - Inches(0.85), Inches(0.35),
             size=12, bold=True, color=TEXT)
    add_text(sl3, desc, rp_l + Inches(0.75), oy + Inches(0.5), rp_w - Inches(0.85), Inches(0.42),
             size=10, color=SUB)
    oy += Inches(1.2)

# 전체 흐름 요약 박스
add_rounded_rect(sl3, rp_l, oy, rp_w, Inches(1.35), fill_color=CARD, line_color=BORDER, radius=0.05)
add_text(sl3, "📌  전체 흐름", rp_l + Inches(0.15), oy + Inches(0.1), rp_w - Inches(0.2), Inches(0.28),
         size=10, bold=True, color=ORANGE)
add_text_multiline(sl3,
    [("DW 매출 자료 수령  →  보정매출 금액 확인 (A/B/C)", False, BODY, 10),
     ("→  매출조정 수기테이블 기입  →  11일 배치 대기", False, BODY, 10)],
    rp_l + Inches(0.15), oy + Inches(0.38), rp_w - Inches(0.2), Inches(0.8))


# ══════════════════════════════════════════════
# SLIDE 4 — 11일 배치 (자동)
# ══════════════════════════════════════════════
sl4 = prs.slides.add_slide(BLANK)
slide_bg(sl4, BG)

# 헤더 (어두운 톤)
hbar = add_rect(sl4, 0, 0, W, Inches(1.1), fill_color=RGBColor(0x3D,0x3D,0x3A))
add_text(sl4, "자동 실행 · 담당자 액션 없음", Inches(0.4), Inches(0.08), Inches(10), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(sl4, "매월 11일 — 데이터 배치 자동화", Inches(0.4), Inches(0.3), Inches(12), Inches(0.55),
         size=22, bold=True, color=WHITE)
add_text(sl4, "스케줄러가 DW 원장 생성 감지 시점에 후속 배치 자동 트리거", Inches(0.4), Inches(0.82),
         Inches(12), Inches(0.3), size=11, color=RGBColor(0xFF,0xCC,0xAA))

# 왼쪽: 세로 타임라인
tl_items = [
    ("01", MUTED,   BORDER,  BODY,    "매월 6 영업일",      "기업 DW 확정매출 및 부문별 PL 계정 산출 마감\n달력상 9~10일경 유동"),
    ("02", CARD,    ORANGE,  ORANGE,  "매월 11일 04:00",    "DW → 분석 데이터레이크(DL) 이관 및 적재 자동 완료\n공휴일 무관"),
    ("03", ORANGE,  None,    WHITE,   "매월 11일 05:00~",   "DL 원장 조인 → 성과관리 데이터 마트 구축\n스케줄러 감지 후 연쇄 기동"),
]

tl_l = Inches(0.5); tl_t = Inches(1.4)
dot_d = Inches(0.56); body_x = tl_l + dot_d + Inches(0.2)
for i, (num, df, dl, tc, title, desc) in enumerate(tl_items):
    y = tl_t + i * Inches(1.5)
    s = sl4.shapes.add_shape(5, tl_l, y, dot_d, dot_d)
    s.adjustments[0] = 0.1
    s.fill.solid(); s.fill.fore_color.rgb = df
    if dl:
        s.line.color.rgb = dl; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    add_text(sl4, num, tl_l, y, dot_d, dot_d, size=14, bold=True, color=tc, align=PP_ALIGN.CENTER)

    add_text(sl4, title, body_x, y, Inches(5.5), Inches(0.4), size=15, bold=True, color=TEXT)
    add_text(sl4, desc,  body_x, y + Inches(0.42), Inches(5.5), Inches(0.9), size=10, color=BODY)

    if i < len(tl_items) - 1:
        add_rect(sl4, tl_l + dot_d/2 - Inches(0.01), y + dot_d,
                  Inches(0.02), Inches(1.5) - dot_d, fill_color=BORDER)

# 자동화 안내 박스
auto_y = tl_t + len(tl_items) * Inches(1.5) + Inches(0.1)
add_rounded_rect(sl4, Inches(0.5), auto_y, Inches(6.2), Inches(0.65),
                  fill_color=GREEN_BG, line_color=RGBColor(0xA8,0xD9,0xD0), radius=0.06)
add_text_multiline(sl4,
    [("✓  연휴·마감 지연 시에도 ", False, GREEN, 11),],
    Inches(0.65), auto_y + Inches(0.08), Inches(5.8), Inches(0.28))
add_text(sl4, "인위적 수동 조정 불필요 — 스케줄러가 DW 원장 생성 감지 즉시 자동 트리거",
         Inches(0.65), auto_y + Inches(0.32), Inches(5.8), Inches(0.28), size=10, color=GREEN)

# 오른쪽: 분기 안내
rp_x = Inches(7.0); rp_t2 = Inches(1.4); rp_w2 = Inches(5.9)

add_text(sl4, "배치 완료 후 분기 확인", rp_x, rp_t2, rp_w2, Inches(0.32),
         size=13, bold=True, color=TEXT)

# IF 박스
add_rounded_rect(sl4, rp_x, rp_t2 + Inches(0.42), rp_w2, Inches(2.1),
                  fill_color=RGBColor(0xF0,0xEE,0xE5), line_color=BORDER, radius=0.06)
pill_badge(sl4, rp_x + Inches(0.15), rp_t2 + Inches(0.55), "IF", fill=DARK, text_color=WHITE, size=9)
add_text(sl4, "수기테이블 미완성 항목 있을 시",
         rp_x + Inches(0.6), rp_t2 + Inches(0.52), rp_w2 - Inches(0.7), Inches(0.32),
         size=12, bold=True, color=TEXT)
add_text_multiline(sl4,
    [("1. 누락 수기테이블 항목 완성", False, BODY, 11),
     ("2. 성과관리 데이터 재생성 요청", False, BODY, 11),
     ("3. 완료 후 → ELSE 루틴 수행", False, ORANGE, 11)],
    rp_x + Inches(0.2), rp_t2 + Inches(0.92), rp_w2 - Inches(0.3), Inches(1.0))

# ELSE 박스
else_y = rp_t2 + Inches(2.7)
add_rounded_rect(sl4, rp_x, else_y, rp_w2, Inches(2.1),
                  fill_color=RGBColor(0xFD,0xF0,0xE8), line_color=ORANGE, radius=0.06)
pill_badge(sl4, rp_x + Inches(0.15), else_y + Inches(0.13), "ELSE", fill=ORANGE, text_color=WHITE, size=9)
add_text(sl4, "수기테이블 완성 상태인 경우",
         rp_x + Inches(0.75), else_y + Inches(0.1), rp_w2 - Inches(0.85), Inches(0.32),
         size=12, bold=True, color=TEXT)
add_text_multiline(sl4,
    [("1. KA 고객 매출·이익 및 심의내역 추출", False, BODY, 11),
     ("2. 네이버클라우드·한전 수익성 조정", False, BODY, 11),
     ("3. 나머지 KA 고객 특이사항 Eye-Checking", False, BODY, 11)],
    rp_x + Inches(0.2), else_y + Inches(0.5), rp_w2 - Inches(0.3), Inches(1.4))


# ══════════════════════════════════════════════
# SLIDE 5 — IF: 수기테이블 미완성 시
# ══════════════════════════════════════════════
sl5 = prs.slides.add_slide(BLANK)
slide_bg(sl5, BG)

# 헤더 (어두운)
add_rect(sl5, 0, 0, W, Inches(1.1), fill_color=DARK)
add_text(sl5, "조건부 실행", Inches(0.4), Inches(0.08), Inches(5), Inches(0.28),
         size=9, bold=True, color=ORANGE)
add_text(sl5, "IF — 수기테이블 미완성 시 대응",
         Inches(0.4), Inches(0.28), Inches(12), Inches(0.55), size=22, bold=True, color=WHITE)
add_text(sl5, "배치 실행 후 누락 항목이 있을 경우",
         Inches(0.4), Inches(0.82), Inches(12), Inches(0.28), size=11, color=RGBColor(0xFF,0xCC,0xAA))

# IF 핵심 박스
add_rounded_rect(sl5, Inches(0.5), Inches(1.3), Inches(8.8), Inches(2.1),
                  fill_color=DARK, line_color=None, radius=0.05)
add_text(sl5, "IF: 수기테이블 미완성", Inches(0.7), Inches(1.38), Inches(7), Inches(0.28),
         size=9, bold=True, color=ORANGE)
add_text(sl5, "데이터 재생성 요청", Inches(0.7), Inches(1.65), Inches(7), Inches(0.5),
         size=20, bold=True, color=WHITE)
add_text_multiline(sl5,
    [("배치 완료 후 수기테이블 누락 항목을 추가로 기입한 뒤,", False, RGBColor(0xFF,0xE0,0xCC), 11),
     ("담당자에게 성과관리 데이터 재생성을 요청합니다.", False, RGBColor(0xFF,0xE0,0xCC), 11)],
    Inches(0.7), Inches(2.18), Inches(7.8), Inches(0.7))

# 화살표 힌트
add_rounded_rect(sl5, Inches(9.5), Inches(1.5), Inches(3.4), Inches(0.65),
                  fill_color=ORANGE, radius=0.12)
add_text(sl5, "재생성 완료 →  ELSE 루틴 수행 (슬라이드 6)",
         Inches(9.55), Inches(1.55), Inches(3.3), Inches(0.52),
         size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 재생성 전 체크리스트
add_text(sl5, "재생성 전 체크", Inches(0.5), Inches(3.6), Inches(6), Inches(0.3),
         size=11, bold=True, color=SUB)

checks5 = [
    ("에스원·KGM·토요타  보정매출 전부 기입 완료 확인", ""),
    ("한전 PL 수익률  수기테이블 최신월 입력 완료 확인", ""),
    ("재생성 완료 후 최신월 마트 적재 시간 확인", ""),
]
cy5 = Inches(4.0)
for txt, sub in checks5:
    rh = check_row(sl5, Inches(0.5), cy5, Inches(8.8), txt, sub)
    cy5 += rh + Inches(0.12)

# 우측: 누락 원인 안내
add_rounded_rect(sl5, Inches(9.5), Inches(2.4), Inches(3.4), Inches(4.6),
                  fill_color=CARD, line_color=BORDER, radius=0.05)
add_text(sl5, "주요 누락 체크 항목", Inches(9.65), Inches(2.52), Inches(3.1), Inches(0.3),
         size=10, bold=True, color=ORANGE)

nurak = [
    ("에스원  M21 / M2M"),
    ("KGM  T35 / 커넥티드카"),
    ("토요타  T35 / 커넥티드카"),
    ("한전 PL 수익률 (P10)"),
]
ny = Inches(2.9)
for item in nurak:
    add_rounded_rect(sl5, Inches(9.62), ny, Inches(3.1), Inches(0.48),
                      fill_color=MUTED, line_color=BORDER, radius=0.05)
    add_text(sl5, item, Inches(9.78), ny + Inches(0.08), Inches(2.85), Inches(0.32),
             size=10, bold=True, color=TEXT)
    ny += Inches(0.6)


# ══════════════════════════════════════════════
# SLIDE 6 — ELSE: 수익성 조정
# ══════════════════════════════════════════════
sl6 = prs.slides.add_slide(BLANK)
slide_bg(sl6, BG)
orange_header(sl6, "수익성 조정 — 네이버클라우드 · 한전",
              "배치 직후 예외 고객사 조정처리 반영", "02 · 데이터 생성 완료 후")

# 왼쪽: 네이버클라우드
col_w = Inches(6.1)
add_rounded_rect(sl6, Inches(0.4), Inches(1.3), col_w, Inches(5.8),
                  fill_color=CARD, line_color=BORDER, radius=0.04)
pill_badge(sl6, Inches(0.55), Inches(1.4), "예외 #1", fill=CARD, text_color=ORANGE, size=9)
add_text(sl6, "네이버클라우드 (코람코 IDC)", Inches(0.55), Inches(1.7), col_w - Inches(0.2), Inches(0.38),
         size=14, bold=True, color=TEXT)
add_text(sl6, "Billing 역산 예외처리 — 3단계", Inches(0.55), Inches(2.1), col_w - Inches(0.2), Inches(0.28),
         size=10, color=SUB)

checks6_n = [
    ("UCube에서 네이버클라우드 청구계정  청구금액 확인", ""),
    ("11일 이후 최신월 성과관리 데이터  생성 완료 확인", ""),
    ("심의내역에서 네이버 IDC 코람코 관련 심의 찾아\n반영매출을 1월~최신월 누적값으로 조정처리 반영", ""),
]
cy6n = Inches(2.5)
for txt, sub in checks6_n:
    rh = check_row(sl6, Inches(0.55), cy6n, col_w - Inches(0.2), txt, sub)
    cy6n += rh + Inches(0.12)

# 왜곡 원인 안내박스
add_rounded_rect(sl6, Inches(0.55), cy6n + Inches(0.05), col_w - Inches(0.2), Inches(1.1),
                  fill_color=RGBColor(0xFF,0xF5,0xEC), line_color=ORANGE, radius=0.05)
add_text(sl6, "✗ 왜곡 원인", Inches(0.7), cy6n + Inches(0.12), Inches(4.5), Inches(0.25),
         size=9, bold=True, color=ORANGE)
add_text(sl6, "심의 이익률 약 1.53%로 극히 낮고 초기 매출 과소 계상 → 이익률 왜곡\n→ 실제 청구 매출(Billing) 기준으로 역산 적용",
         Inches(0.7), cy6n + Inches(0.35), col_w - Inches(0.3), Inches(0.65),
         size=9, color=BODY)

# 오른쪽: 한전
rp6_l = Inches(6.8)
add_rounded_rect(sl6, rp6_l, Inches(1.3), col_w, Inches(5.8),
                  fill_color=CARD, line_color=BORDER, radius=0.04)
pill_badge(sl6, rp6_l + Inches(0.15), Inches(1.4), "예외 #2", fill=MUTED, text_color=SUB, size=9)
add_text(sl6, "한전 KEPCO 전력통신", rp6_l + Inches(0.15), Inches(1.7), col_w - Inches(0.2), Inches(0.38),
         size=14, bold=True, color=TEXT)
add_text(sl6, "PL 이익률 강제 적용 — 4단계", rp6_l + Inches(0.15), Inches(2.1), col_w - Inches(0.2), Inches(0.28),
         size=10, color=SUB)

checks6_k = [
    ("6영업일 즈음 기획팀에서 전달하는  PL 자료 수령", ""),
    ("PL 자료로 PL 수익률 수기테이블 완성\n→ 최신월 P10 전력통신 수익률 입력", ""),
    ("11일 이후 최신월 성과관리 데이터  생성 완료 확인", ""),
    ("심의내역에서 한전 P10 전력통신 심의 찾아\n수익률을 PL 수익률로 조정처리 반영", ""),
]
cy6k = Inches(2.5)
for txt, sub in checks6_k:
    rh = check_row(sl6, rp6_l + Inches(0.15), cy6k, col_w - Inches(0.2), txt, sub)
    cy6k += rh + Inches(0.12)

add_rounded_rect(sl6, rp6_l + Inches(0.15), cy6k + Inches(0.05), col_w - Inches(0.2), Inches(0.7),
                  fill_color=RGBColor(0xFF,0xF5,0xEC), line_color=ORANGE, radius=0.05)
add_text(sl6, "✗ 왜곡 원인", rp6_l + Inches(0.3), cy6k + Inches(0.1), Inches(4), Inches(0.22),
         size=9, bold=True, color=ORANGE)
add_text(sl6, "한전↔전력통신 상호집중 구조 → VRB 심의 극소 → 일반 로직 정합성 훼손",
         rp6_l + Inches(0.3), cy6k + Inches(0.3), col_w - Inches(0.4), Inches(0.3),
         size=9, color=BODY)


# ══════════════════════════════════════════════
# SLIDE 7 — Eye-Checking
# ══════════════════════════════════════════════
sl7 = prs.slides.add_slide(BLANK)
slide_bg(sl7, BG)
orange_header(sl7, "Eye-Checking — KA 고객 특이사항 확인",
              "Claude로 체크 양식 생성 → 항목별 육안 확인", "03 · 최종 확인")

# 절차 (2단계 화살표)
proc_y = Inches(1.3)
for i, (num, title, sub) in enumerate([
    ("1", "Claude에 체크 양식 요청", "KA 고객 목록 + 해당 월 성과 데이터 첨부"),
    ("2", "항목별 육안 확인",        "이상 항목 발견 시 원인 파악 및 수정"),
]):
    px = Inches(0.4 + i * 5.8)
    add_rounded_rect(sl7, px, proc_y, Inches(5.2), Inches(0.9),
                      fill_color=CARD, line_color=BORDER, radius=0.05)
    add_circle(sl7, px + Inches(0.15), proc_y + Inches(0.17), Inches(0.52), ORANGE)
    add_text(sl7, num, px + Inches(0.15), proc_y + Inches(0.17), Inches(0.52), Inches(0.52),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl7, title, px + Inches(0.78), proc_y + Inches(0.1), Inches(4.2), Inches(0.3),
             size=12, bold=True, color=TEXT)
    add_text(sl7, sub,   px + Inches(0.78), proc_y + Inches(0.42), Inches(4.2), Inches(0.36),
             size=9, color=SUB)
    if i == 0:
        add_text(sl7, "→", Inches(5.65), proc_y + Inches(0.28), Inches(0.6), Inches(0.4),
                 size=20, bold=True, color=BORDER, align=PP_ALIGN.CENTER)

# 체크 항목 2열
col_titles = ["수치 이상 체크", "데이터 정합성 체크"]
checks7 = [
    [
        ("매출 급변동 — 전월 대비 ±30% 초과 여부", ""),
        ("이익률 이상 — 업계 평균 대비 극단값 여부", ""),
        ("GRR/NRR 이상 — 100% 초과 또는 극단 하락", ""),
        ("수주잔고 갑작스러운 변동 여부", ""),
    ],
    [
        ("심의내역 누락 — VRB 심의 등록 여부", ""),
        ("네이버클라우드·한전 수익성 조정 완료 확인", ""),
        ("에스원·KGM·토요타 보정매출 반영 확인", ""),
        ("신규 고객 마스터 매핑 오류 없음 확인", ""),
    ],
]
col_xs = [Inches(0.4), Inches(6.9)]
col_w7 = Inches(6.1)
for ci, (ctitle, citems) in enumerate(zip(col_titles, checks7)):
    cx = col_xs[ci]
    add_text(sl7, ctitle, cx, Inches(2.45), col_w7, Inches(0.28),
             size=10, bold=True, color=SUB)
    cy7 = Inches(2.8)
    for txt, sub in citems:
        rh = check_row(sl7, cx, cy7, col_w7, txt, sub)
        cy7 += rh + Inches(0.1)


# ══════════════════════════════════════════════
# SLIDE 8 — 비정기 작업
# ══════════════════════════════════════════════
sl8 = prs.slides.add_slide(BLANK)
slide_bg(sl8, BG)

# 헤더 (회색)
add_rect(sl8, 0, 0, W, Inches(1.1), fill_color=SUB)
add_text(sl8, "비정기 · 필요시만 수행", Inches(0.4), Inches(0.08), Inches(8), Inches(0.28),
         size=9, bold=True, color=RGBColor(0xD0,0xCD,0xC5))
add_text(sl8, "수기테이블 정의 업데이트 · 고객 마스터 정비",
         Inches(0.4), Inches(0.28), Inches(12), Inches(0.55), size=22, bold=True, color=WHITE)
add_text(sl8, "변경 사유 발생 시 즉시 처리 — 정기 주기 없음",
         Inches(0.4), Inches(0.82), Inches(12), Inches(0.28), size=11, color=RGBColor(0xD0,0xCD,0xC5))

# 왼쪽: 수기테이블 정의 업데이트
add_text(sl8, "수기테이블 정의 업데이트", Inches(0.4), Inches(1.28), Inches(5.8), Inches(0.3),
         size=11, bold=True, color=SUB)

irreg_l = [
    ("1", "KA 고객 정의",       "Key Account 식별 기준 변경 시"),
    ("2", "GA 고객 정의",       "General Account 기준 변경 시"),
    ("3", "고객 그룹핑 정의",   "복수 고객 단일 그룹 매핑 변경 시"),
    ("4", "고객군 정의",        "산업군 세그먼트 분류 변경 시"),
]
iy = Inches(1.65)
for num, title, desc in irreg_l:
    add_rounded_rect(sl8, Inches(0.4), iy, Inches(6.0), Inches(0.7),
                      fill_color=MUTED, line_color=BORDER, radius=0.05)
    add_rounded_rect(sl8, Inches(0.55), iy + Inches(0.12), Inches(0.45), Inches(0.45),
                      fill_color=CARD, line_color=BORDER, radius=0.1)
    add_text(sl8, num, Inches(0.55), iy + Inches(0.12), Inches(0.45), Inches(0.45),
             size=11, bold=True, color=SUB, align=PP_ALIGN.CENTER)
    add_text(sl8, title, Inches(1.1), iy + Inches(0.08), Inches(3.8), Inches(0.3),
             size=12, bold=True, color=TEXT)
    add_text(sl8, desc, Inches(1.1), iy + Inches(0.38), Inches(3.8), Inches(0.25),
             size=9, color=SUB)
    # 필요시 뱃지
    add_rounded_rect(sl8, Inches(5.5), iy + Inches(0.18), Inches(0.75), Inches(0.3),
                      fill_color=MUTED, line_color=BORDER, radius=0.3)
    add_text(sl8, "필요시", Inches(5.5), iy + Inches(0.18), Inches(0.75), Inches(0.3),
             size=8, bold=True, color=SUB, align=PP_ALIGN.CENTER)
    iy += Inches(0.82)

# 오른쪽: 고객 마스터 정비
add_text(sl8, "고객 마스터 정비", Inches(6.9), Inches(1.28), Inches(5.8), Inches(0.3),
         size=11, bold=True, color=SUB)

irreg_r = [
    ("🔗", "사업자번호 매핑 오류 수정",  "중복 계정 정제 / 유일 고객 ID 재발급", "오류 발견 시"),
    ("➕", "신규 고객 등록",              "사업자번호 확인 후 마스터 추가",       "신규 발생 시"),
    ("📋", "수기테이블 항목 수 변동",     "현재 18개 수기테이블 구조 변경 시",   "구조 변경 시"),
]
iy2 = Inches(1.65)
for icon, title, desc, tag in irreg_r:
    add_rounded_rect(sl8, Inches(6.9), iy2, Inches(6.0), Inches(0.7),
                      fill_color=MUTED, line_color=BORDER, radius=0.05)
    add_rounded_rect(sl8, Inches(7.05), iy2 + Inches(0.12), Inches(0.45), Inches(0.45),
                      fill_color=CARD, line_color=BORDER, radius=0.1)
    add_text(sl8, icon, Inches(7.05), iy2 + Inches(0.12), Inches(0.45), Inches(0.45),
             size=14, align=PP_ALIGN.CENTER, color=TEXT)
    add_text(sl8, title, Inches(7.62), iy2 + Inches(0.08), Inches(4.0), Inches(0.3),
             size=12, bold=True, color=TEXT)
    add_text(sl8, desc, Inches(7.62), iy2 + Inches(0.38), Inches(4.0), Inches(0.25),
             size=9, color=SUB)
    bw = Inches(len(tag) * 0.12 + 0.3)
    add_rounded_rect(sl8, Inches(11.95) - bw, iy2 + Inches(0.18), bw, Inches(0.3),
                      fill_color=MUTED, line_color=BORDER, radius=0.3)
    add_text(sl8, tag, Inches(11.95) - bw, iy2 + Inches(0.18), bw, Inches(0.3),
             size=8, bold=True, color=SUB, align=PP_ALIGN.CENTER)
    iy2 += Inches(0.82)

# 수기테이블 현황 안내 박스
mh_y = iy2 + Inches(0.1)
add_rounded_rect(sl8, Inches(6.9), mh_y, Inches(6.0), Inches(0.85),
                  fill_color=CARD, line_color=BORDER, radius=0.05)
add_text(sl8, "📌  수기테이블 현황", Inches(7.05), mh_y + Inches(0.08), Inches(5.7), Inches(0.28),
         size=10, bold=True, color=ORANGE)
add_text_multiline(sl8,
    [("전체 18개 중  주요 관리 7개만 매월 점검", False, BODY, 10),
     ("나머지 11개는 필요시 담당자 확인", False, SUB, 9)],
    Inches(7.05), mh_y + Inches(0.35), Inches(5.6), Inches(0.42))

# 마무리 박스
fin_y = Inches(6.55)
add_rounded_rect(sl8, Inches(0.4), fin_y, Inches(12.5), Inches(0.7),
                  fill_color=ORANGE, radius=0.06)
add_text(sl8, "이번 달 작업 완료 🎉   다음 달 11일 배치 전까지 수기테이블 수령 즉시 대응",
         Inches(0.6), fin_y + Inches(0.15), Inches(12.1), Inches(0.4),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── 저장 ──────────────────────────────────────
out_path = "/Users/konslie/Desktop/ClaudeCode/인수인계쨩/성과관리_월별작업_체크리스트.pptx"
prs.save(out_path)
print(f"저장 완료: {out_path}")
